import os
import re
from pptx import Presentation

# 解析LRC字幕文件
def parse_lrc(lrc_file):
    """解析LRC字幕文件，返回时间点和对应文本的列表"""
    lrc_data = []
    
    with open(lrc_file, 'r', encoding='utf-8') as f:
        lines = f.readlines()
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
        
        # 匹配LRC时间格式 [mm:ss.xx]
        time_matches = re.findall(r'\[(\d+):(\d+\.\d+)\]', line)
        if time_matches:
            text = re.sub(r'\[\d+:\d+\.\d+\]', '', line).strip()
            for match in time_matches:
                minutes, seconds = match
                total_seconds = int(minutes) * 60 + float(seconds)
                lrc_data.append((total_seconds, text))
    
    # 按时间排序
    lrc_data.sort(key=lambda x: x[0])
    return lrc_data

# 读取PPT内容
def read_ppt_content(ppt_file):
    """读取PPT每页的文本内容"""
    prs = Presentation(ppt_file)
    slides_content = []
    
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                slide_text.append(shape.text.strip())
        # 合并每页的文本
        page_content = ' '.join(slide_text)
        slides_content.append(page_content)
    
    return slides_content

# 计算每页PPT的放映时间
def calculate_slide_timings(slides_content, lrc_data):
    """根据LRC字幕计算每页PPT的放映时间"""
    if not slides_content or not lrc_data:
        return []
    
    timings = []
    current_slide = 0
    start_time = lrc_data[0][0]
    
    for i, (time, text) in enumerate(lrc_data):
        # 检查当前文本是否与当前幻灯片内容相关
        # 简单匹配：检查文本是否在幻灯片内容中
        if current_slide < len(slides_content):
            slide_content = slides_content[current_slide]
            
            # 如果文本不在当前幻灯片内容中，可能需要切换幻灯片
            if text and text not in slide_content:
                # 计算当前幻灯片的放映时间
                duration = time - start_time
                timings.append(duration)
                
                # 移动到下一张幻灯片
                current_slide += 1
                start_time = time
    
    # 处理最后一张幻灯片
    if current_slide < len(slides_content):
        last_time = lrc_data[-1][0]
        duration = last_time - start_time
        timings.append(duration)
    
    # 确保所有幻灯片都有时间
    while len(timings) < len(slides_content):
        timings.append(5.0)  # 默认5秒
    
    return timings

# 设置PPT每页的放映时间
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
def set_ppt_timings(ppt_file, timings):
    """设置PPT每页的放映时间（完全手动操作XML，兼容所有版本）"""
    prs = Presentation(ppt_file)
    
    for i, slide in enumerate(prs.slides):
        if i < len(timings):
            # 获取幻灯片的根元素
            slide_elem = slide.element
            
            # 查找现有的 <p:transition> 元素
            transition = slide_elem.find(qn('p:transition'))
            if transition is None:
                # 创建新的 <p:transition> 元素
                transition = parse_xml(
                    '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>'
                )
                slide_elem.append(transition)
            
            # 设置自动切换属性
            transition.set('advanceOnTime', '1')          # 启用按时间切换
            transition.set('advTm', str(int(timings[i] * 1000)))  # 毫秒单位
                
    output_ppt = os.path.splitext(ppt_file)[0] + '_timed.pptx'
    prs.save(output_ppt)
    return output_ppt
# def set_ppt_timings(ppt_file, timings):
#     """设置PPT每页的放映时间"""
#     prs = Presentation(ppt_file)
    
#     for i, slide in enumerate(prs.slides):
#         if i < len(timings):
#             # 正确设置幻灯片切换时间
#             slide.slide_transition.advance_on_time = True
#             slide.slide_transition.advance_time = timings[i]
    
#     # 保存修改后的PPT
#     output_ppt = os.path.splitext(ppt_file)[0] + '_timed.pptx'
#     prs.save(output_ppt)
#     return output_ppt

# 主函数
def main():
    # 文件路径
    lrc_file = 'build\\audio_files\\textCleaned_tongyi-xiaomi-analysis-flash.lrc'
    text_file = 'textCleaned_tongyi-xiaomi-analysis-flash.txt'
    ppt_file = '望君长留kp团.pptx'
    
    # 检查文件是否存在
    if not os.path.exists(lrc_file):
        print(f"错误：字幕文件 {lrc_file} 不存在")
        return
    
    if not os.path.exists(ppt_file):
        print(f"错误：PPT文件 {ppt_file} 不存在")
        return
    
    # 解析LRC字幕
    print("解析字幕文件...")
    lrc_data = parse_lrc(lrc_file)
    
    # 读取PPT内容
    print("读取PPT内容...")
    slides_content = read_ppt_content(ppt_file)
    
    # 计算每页放映时间
    print("计算放映时间...")
    timings = calculate_slide_timings(slides_content, lrc_data)
    
    # 设置PPT放映时间
    print("设置PPT放映时间...")
    output_ppt = set_ppt_timings(ppt_file, timings)
    
    print(f"处理完成！\n输出文件：{output_ppt}")
    print(f"每页放映时间：{timings}")

if __name__ == "__main__":
    main()
