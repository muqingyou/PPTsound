import sys
import os
import json
import requests
import time
import hmac
import hashlib
import base64
from pptx import Presentation
from pptx.util import Inches
from PyQt5.QtWidgets import (
    QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout,
    QLabel, QLineEdit, QPushButton, QFileDialog, QComboBox,
    QSpinBox, QCheckBox, QTextEdit, QTabWidget, QListWidget,
    QListWidgetItem, QGroupBox, QMessageBox
)
from PyQt5.QtCore import Qt, QThread, pyqtSignal

# 全局配置
CONFIG = {
    "ALIYUN_ACCESS_KEY_ID": "",
    "ALIYUN_ACCESS_KEY_SECRET": "",
    "ALIYUN_APP_KEY": "",
    "ALIYUN_REGION_ID": "cn-shanghai"
}

# 人物配置
CHARACTERS = []

# 音色列表
VOICES = []

class TestVoiceThread(QThread):
    """测试音色的线程类"""
    signal = pyqtSignal(str, str)
    
    def __init__(self, voice_name, test_text, test_audio_file):
        super().__init__()
        self.voice_name = voice_name
        self.test_text = test_text
        self.test_audio_file = test_audio_file
    
    def run(self):
        try:
            # 获取token
            def get_token():
                url = f"https://nls-meta.{CONFIG['ALIYUN_REGION_ID']}.aliyuncs.com/?Action=CreateToken&Version=2019-02-28"
                timestamp = time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())
                signature_nonce = str(int(time.time() * 1000))
                
                import urllib.parse
                params = {
                    "AccessKeyId": CONFIG["ALIYUN_ACCESS_KEY_ID"],
                    "Action": "CreateToken",
                    "Format": "JSON",
                    "SignatureMethod": "HMAC-SHA1",
                    "SignatureNonce": signature_nonce,
                    "SignatureVersion": "1.0",
                    "Timestamp": timestamp,
                    "Version": "2019-02-28"
                }
                
                sorted_params = sorted(params.items(), key=lambda x: x[0])
                canonical_query_string = "&".join([f"{k}={urllib.parse.quote(v, safe='')}" for k, v in sorted_params])
                string_to_sign = f"GET&%2F&{urllib.parse.quote(canonical_query_string, safe='')}"
                
                signature = base64.b64encode(hmac.new(
                    (CONFIG["ALIYUN_ACCESS_KEY_SECRET"] + "&").encode(),
                    string_to_sign.encode(),
                    hashlib.sha1
                ).digest()).decode()
                
                full_url = f"{url}&AccessKeyId={CONFIG['ALIYUN_ACCESS_KEY_ID']}&Signature={urllib.parse.quote(signature, safe='')}&SignatureNonce={signature_nonce}&Timestamp={urllib.parse.quote(timestamp, safe='')}&Format=JSON&SignatureMethod=HMAC-SHA1&SignatureVersion=1.0"
                
                response = requests.get(full_url)
                if response.status_code == 200:
                    result = response.json()
                    if result.get("Token"):
                        return result["Token"]["Id"]
                return None
            
            token = get_token()
            if not token:
                self.signal.emit("error", "获取TTS token失败")
                return
            
            # 调用TTS API
            url = f"https://nls-gateway.{CONFIG['ALIYUN_REGION_ID']}.aliyuncs.com/stream/v1/tts"
            headers = {
                "Content-Type": "application/json",
                "X-NLS-Token": token
            }
            
            payload = {
                "appkey": CONFIG["ALIYUN_APP_KEY"],
                "text": self.test_text,
                "voice": self.voice_name,
                "format": "mp3",
                "sample_rate": 16000
            }
            
            response = requests.post(url, headers=headers, json=payload)
            if response.status_code == 200:
                with open(self.test_audio_file, "wb") as f:
                    f.write(response.content)
                
                # 尝试播放音频
                try:
                    import platform
                    system = platform.system()
                    if system == "Windows":
                        os.startfile(self.test_audio_file)
                    elif system == "Darwin":  # macOS
                        os.system(f"open {self.test_audio_file}")
                    else:  # Linux
                        os.system(f"xdg-open {self.test_audio_file}")
                    self.signal.emit("success", f"音色测试成功，正在播放音频")
                except Exception as e:
                    self.signal.emit("info", f"音色测试成功，但无法自动播放音频: {str(e)}")
                    self.signal.emit("info", f"音频文件保存为: {self.test_audio_file}")
            else:
                self.signal.emit("error", f"音色测试失败: {response.text}")
        except Exception as e:
            self.signal.emit("error", f"测试音色时发生错误: {str(e)}")

class WorkerThread(QThread):
    """工作线程，用于处理耗时操作"""
    signal = pyqtSignal(str, str)  # 信号类型：(消息类型, 消息内容)
    
    def __init__(self, task_type, **kwargs):
        super().__init__()
        self.task_type = task_type
        self.kwargs = kwargs
    
    def run(self):
        if self.task_type == "get_voices":
            self.get_aliyun_voices()
        elif self.task_type == "process_ppt":
            self.process_ppt()
    
    def get_aliyun_voices(self):
        """获取阿里云支持的音色"""
        try:
            # 从音色.xlsx文件加载音色列表
            voice_file = "音色.xlsx"
            if os.path.exists(voice_file):
                try:
                    import pandas as pd
                    
                    # 读取Excel文件
                    df = pd.read_excel(voice_file)
                    
                    # 解析音色列表
                    voices = []
                    for i, row in df.iterrows():
                        # 尝试获取音色名称、性别和描述
                        voice_name = str(row.iloc[0]).strip() if not pd.isna(row.iloc[0]) else f"voice_{i+1}"
                        
                        
                        # 尝试获取描述
                        description = f"音色 {i+1}"
                        if len(row) > 1 and not pd.isna(row.iloc[1]):
                            description = str(row.iloc[1]).strip()
                        
                        voices.append({
                            "name": voice_name,
                            "description": description
                        })
                    
                    global VOICES
                    VOICES = voices
                    self.signal.emit("success", f"从Excel文件加载音色列表成功，共 {len(voices)} 个音色")
                except Exception as e:
                    self.signal.emit("error", f"读取Excel文件失败: {str(e)}")
        except Exception as e:
            self.signal.emit("error", f"获取音色列表失败: {str(e)}")
        
    def get_aliyun_tts_token(self):
        """获取阿里云TTS服务的访问令牌"""
        url = f"https://nls-meta.{CONFIG['ALIYUN_REGION_ID']}.aliyuncs.com/?Action=CreateToken&Version=2019-02-28"
        
        # 生成签名（使用ISO 8601格式的时间戳）
        timestamp = time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())
        signature_nonce = str(int(time.time() * 1000))
        
        # 构建签名字符串（需要对参数值进行URL编码）
        import urllib.parse
        params = {
            "AccessKeyId": CONFIG["ALIYUN_ACCESS_KEY_ID"],
            "Action": "CreateToken",
            "Format": "JSON",
            "SignatureMethod": "HMAC-SHA1",
            "SignatureNonce": signature_nonce,
            "SignatureVersion": "1.0",
            "Timestamp": timestamp,
            "Version": "2019-02-28"
        }
        
        # 对参数进行排序并构建签名字符串
        sorted_params = sorted(params.items(), key=lambda x: x[0])
        canonical_query_string = "&".join([f"{k}={urllib.parse.quote(v, safe='')}" for k, v in sorted_params])
        string_to_sign = f"GET&%2F&{urllib.parse.quote(canonical_query_string, safe='')}"
        
        # 计算签名
        signature = base64.b64encode(hmac.new(
            (CONFIG["ALIYUN_ACCESS_KEY_SECRET"] + "&").encode(),
            string_to_sign.encode(),
            hashlib.sha1
        ).digest()).decode()
        
        # 构建完整URL
        full_url = f"{url}&AccessKeyId={CONFIG['ALIYUN_ACCESS_KEY_ID']}&Signature={urllib.parse.quote(signature, safe='')}&SignatureNonce={signature_nonce}&Timestamp={urllib.parse.quote(timestamp, safe='')}&Format=JSON&SignatureMethod=HMAC-SHA1&SignatureVersion=1.0"
        
        response = requests.get(full_url)
        if response.status_code == 200:
            result = response.json()
            if result.get("Token"):
                return result["Token"]["Id"]
        return None
    
    def text_to_speech(self, text, output_file, voice="zhitian"):
        """使用阿里云TTS将文本转换为语音"""
        token = self.get_aliyun_tts_token()
        if not token:
            return False
        
        url = f"https://nls-gateway.{CONFIG['ALIYUN_REGION_ID']}.aliyuncs.com/stream/v1/tts"
        headers = {
            "Content-Type": "application/json",
            "X-NLS-Token": token
        }
        
        payload = {
            "appkey": CONFIG["ALIYUN_APP_KEY"],
            "text": text,
            "voice": voice,
            "format": "mp3",
            "sample_rate": 16000
        }
        
        response = requests.post(url, headers=headers, json=payload)
        if response.status_code == 200:
            with open(output_file, "wb") as f:
                f.write(response.content)
            return True
        return False
    
    def process_ppt(self):
        """处理PPT文件"""
        ppt_file = self.kwargs.get("ppt_file")
        start_page = self.kwargs.get("start_page")
        end_page = self.kwargs.get("end_page")
        
        if not ppt_file or not os.path.exists(ppt_file):
            self.signal.emit("error", "PPT文件不存在")
            return
        
        try:
            # 提取PPT文本和识别人物
            prs = Presentation(ppt_file)
            slides_content = []
            
            for i, slide in enumerate(prs.slides):
                page_num = i + 1
                if page_num < start_page or page_num > end_page:
                    continue
                
                slide_text = []
                detected_character = None
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text
                        slide_text.append(text)
                        
                        # 检测文本中的人物名称
                        if detected_character is None:
                            for char_info in CHARACTERS:
                                character_name = char_info.get("name")
                                if character_name in text:
                                    detected_character = character_name
                                    break
                
                # 如果没有检测到人物，使用默认
                if not detected_character:
                    detected_character = "KP"
                
                slides_content.append((page_num, "\n".join(slide_text), detected_character))
            
            # 创建音频文件保存目录
            audio_dir = "audio_files"
            if not os.path.exists(audio_dir):
                os.makedirs(audio_dir)
            
            # 生成音频文件
            audio_files = []
            for i, (page_num, text, character) in enumerate(slides_content):
                audio_file = os.path.join(audio_dir, f"slide_{page_num}_{character}.mp3")
                
                # 检查音频文件是否已存在
                if os.path.exists(audio_file):
                    self.signal.emit("info", f"第 {page_num} 页({character})的音频文件已存在，直接使用...")
                    audio_files.append(audio_file)
                elif text.strip():  # 跳过空文本
                    # 获取对应人物的音色
                    voice = "zhitian"  # 默认音色
                    for char_info in CHARACTERS:
                        if char_info.get("name") == character:
                            voice = char_info.get("voice", "zhitian")
                            break
                    self.signal.emit("info", f"正在为第 {page_num} 页({character})生成音频，使用音色: {voice}...")
                    if self.text_to_speech(text, audio_file, voice=voice):
                        self.signal.emit("success", f"第 {page_num} 页({character})的音频生成成功")
                        audio_files.append(audio_file)
                    else:
                        self.signal.emit("error", f"第 {page_num} 页({character})的音频生成失败")
                        audio_files.append(None)
                else:
                    self.signal.emit("info", f"第 {page_num} 页无文本，跳过音频生成")
                    audio_files.append(None)
            
            # 插入音频到PPT
            if audio_files:
                self.signal.emit("info", "正在插入音频到PPT...")
                output_ppt = ppt_file.replace(".pptx", "_with_audio.pptx")
                prs = Presentation(ppt_file)
                
                for i, (page_num, text, character) in enumerate(slides_content):
                    if i < len(audio_files):
                        audio_file = audio_files[i]
                        if audio_file:
                            slide_index = page_num - 1
                            if 0 <= slide_index < len(prs.slides):
                                slide = prs.slides[slide_index]
                                
                                # 添加音频文件到PPT
                                left = Inches(0)
                                top = Inches(0)
                                width = Inches(1)
                                height = Inches(1)  
                                
                                # 在幻灯片上添加音频
                                try:
                                    audio_shape = slide.shapes.add_movie(
                                        audio_file,
                                        left, top, width, height,
                                        poster_frame_image=None,
                                        mime_type='audio/mp3'
                                    )
                                    
                                    self.signal.emit("success", f"已添加音频到第 {page_num} 页")
                                except Exception as e:
                                    self.signal.emit("error", f"添加音频到第 {page_num} 页失败: {str(e)}")
                
                prs.save(output_ppt)
                self.signal.emit("success", f"处理完成，输出文件: {output_ppt}")
            else:
                self.signal.emit("info", "没有生成音频文件")
                
        except Exception as e:
            self.signal.emit("error", f"处理PPT失败: {str(e)}")

class MainWindow(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("PPT语音添加工具")
        self.setGeometry(100, 100, 800, 600)
        
        # 中心部件
        central_widget = QWidget()
        self.setCentralWidget(central_widget)
        
        # 主布局
        main_layout = QVBoxLayout(central_widget)
        
        # 标签页
        self.tab_widget = QTabWidget()
        main_layout.addWidget(self.tab_widget)
        
        # 配置标签页
        self.config_tab = QWidget()
        self.tab_widget.addTab(self.config_tab, "阿里云配置")
        
        # 人物标签页
        self.character_tab = QWidget()
        self.tab_widget.addTab(self.character_tab, "人物配置")
        
        # PPT处理标签页
        self.ppt_tab = QWidget()
        self.tab_widget.addTab(self.ppt_tab, "PPT处理")
        
        # 初始化配置标签页
        self.init_config_tab()
        
        # 初始化人物标签页
        self.init_character_tab()
        
        # 初始化PPT处理标签页
        self.init_ppt_tab()
        
        # 状态输出
        self.status_text = QTextEdit()
        self.status_text.setReadOnly(True)
        self.status_text.setMinimumHeight(100)
        main_layout.addWidget(QLabel("操作状态:"))
        main_layout.addWidget(self.status_text)
    
    def init_config_tab(self):
        """初始化配置标签页"""
        layout = QVBoxLayout(self.config_tab)
        
        # 配置输入组
        config_group = QGroupBox("阿里云配置")
        config_layout = QVBoxLayout(config_group)
        
        # Access Key ID
        hbox1 = QHBoxLayout()
        hbox1.addWidget(QLabel("Access Key ID:"))
        self.access_key_id = QLineEdit()
        hbox1.addWidget(self.access_key_id)
        config_layout.addLayout(hbox1)
        
        # Access Key Secret
        hbox2 = QHBoxLayout()
        hbox2.addWidget(QLabel("Access Key Secret:"))
        self.access_key_secret = QLineEdit()
        self.access_key_secret.setEchoMode(QLineEdit.Password)
        hbox2.addWidget(self.access_key_secret)
        config_layout.addLayout(hbox2)
        
        # App Key
        hbox3 = QHBoxLayout()
        hbox3.addWidget(QLabel("App Key:"))
        self.app_key = QLineEdit()
        hbox3.addWidget(self.app_key)
        config_layout.addLayout(hbox3)
        
        # 区域ID
        hbox4 = QHBoxLayout()
        hbox4.addWidget(QLabel("区域ID:"))
        self.region_id = QComboBox()
        self.region_id.addItems(["cn-shanghai", "cn-beijing", "cn-hangzhou"])
        self.region_id.setCurrentText("cn-shanghai")
        hbox4.addWidget(self.region_id)
        config_layout.addLayout(hbox4)
        
        layout.addWidget(config_group)
        
        # 配置文件
        file_group = QGroupBox("配置文件")
        file_layout = QHBoxLayout(file_group)
        self.config_file_path = QLineEdit()
        file_layout.addWidget(self.config_file_path)
        browse_btn = QPushButton("浏览")
        browse_btn.clicked.connect(self.browse_config_file)
        file_layout.addWidget(browse_btn)
        load_btn = QPushButton("加载配置")
        load_btn.clicked.connect(self.load_config_file)
        file_layout.addWidget(load_btn)
        layout.addWidget(file_group)
        
        # 保存配置
        save_btn = QPushButton("保存配置")
        save_btn.clicked.connect(self.save_config)
        layout.addWidget(save_btn)
        
        # 测试连接
        test_btn = QPushButton("测试连接")
        test_btn.clicked.connect(self.test_connection)
        layout.addWidget(test_btn)
    
    def init_character_tab(self):
        """初始化人物标签页"""
        layout = QVBoxLayout(self.character_tab)
        
        # 人物管理
        char_group = QGroupBox("人物管理")
        char_layout = QVBoxLayout(char_group)
        
        # 人物输入
        hbox1 = QHBoxLayout()
        hbox1.addWidget(QLabel("人物名称:"))
        self.character_name = QLineEdit()
        hbox1.addWidget(self.character_name)
        char_layout.addLayout(hbox1)
        
        # 音色选择
        hbox2 = QHBoxLayout()
        hbox2.addWidget(QLabel("音色:"))
        self.character_voice = QComboBox()
        hbox2.addWidget(self.character_voice)
        char_layout.addLayout(hbox2)
        
        # 添加/删除人物
        hbox3 = QHBoxLayout()
        add_char_btn = QPushButton("添加人物")
        add_char_btn.clicked.connect(self.add_character)
        hbox3.addWidget(add_char_btn)
        delete_char_btn = QPushButton("删除选中人物")
        delete_char_btn.clicked.connect(self.delete_character)
        hbox3.addWidget(delete_char_btn)
        char_layout.addLayout(hbox3)
        
        # 人物列表
        self.character_list = QListWidget()
        char_layout.addWidget(QLabel("人物列表:"))
        char_layout.addWidget(self.character_list)
        layout.addWidget(char_group)
        
        # 音色管理
        voice_group = QGroupBox("音色管理")
        voice_layout = QVBoxLayout(voice_group)
        
        # 获取音色列表
        get_voices_btn = QPushButton("获取音色列表")
        get_voices_btn.clicked.connect(self.get_voices)
        voice_layout.addWidget(get_voices_btn)
        
        # 测试音色
        test_voice_btn = QPushButton("测试选中音色")
        test_voice_btn.clicked.connect(self.test_selected_voice)
        voice_layout.addWidget(test_voice_btn)
        
        # 音色列表
        self.voice_list = QListWidget()
        voice_layout.addWidget(QLabel("阿里云音色列表:"))
        voice_layout.addWidget(self.voice_list)
        layout.addWidget(voice_group)
    
    def init_ppt_tab(self):
        """初始化PPT处理标签页"""
        layout = QVBoxLayout(self.ppt_tab)
        
        # PPT文件
        hbox1 = QHBoxLayout()
        hbox1.addWidget(QLabel("PPT文件:"))
        self.ppt_file_path = QLineEdit()
        hbox1.addWidget(self.ppt_file_path)
        browse_btn = QPushButton("浏览")
        browse_btn.clicked.connect(self.browse_ppt_file)
        hbox1.addWidget(browse_btn)
        layout.addLayout(hbox1)
        
        # 页面范围
        hbox2 = QHBoxLayout()
        hbox2.addWidget(QLabel("页面范围:"))
        self.start_page = QSpinBox()
        self.start_page.setMinimum(1)
        self.start_page.setValue(1)
        hbox2.addWidget(self.start_page)
        hbox2.addWidget(QLabel("到"))
        self.end_page = QSpinBox()
        self.end_page.setMinimum(1)
        self.end_page.setValue(4)
        hbox2.addWidget(self.end_page)
        layout.addLayout(hbox2)
        
        # 处理按钮
        process_btn = QPushButton("开始处理")
        process_btn.clicked.connect(self.process_ppt)
        layout.addWidget(process_btn)
    
    def browse_config_file(self):
        """浏览配置文件"""
        file_path, _ = QFileDialog.getOpenFileName(self, "选择配置文件", "", "JSON Files (*.json)")
        if file_path:
            self.config_file_path.setText(file_path)
    
    def load_config_file(self):
        """加载配置文件"""
        file_path = self.config_file_path.text()
        if not file_path or not os.path.exists(file_path):
            QMessageBox.warning(self, "警告", "配置文件不存在")
            return
        
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                config = json.load(f)
            
            # 填充配置
            self.access_key_id.setText(config.get("ALIYUN_ACCESS_KEY_ID", ""))
            self.access_key_secret.setText(config.get("ALIYUN_ACCESS_KEY_SECRET", ""))
            self.app_key.setText(config.get("ALIYUN_APP_KEY", ""))
            region_id = config.get("ALIYUN_REGION_ID", "cn-shanghai")
            if region_id in ["cn-shanghai", "cn-beijing", "cn-hangzhou"]:
                self.region_id.setCurrentText(region_id)
            
            # 加载人物配置
            if "CHARACTERS" in config:
                global CHARACTERS
                CHARACTERS = config["CHARACTERS"]
                self.update_character_list()
                self.log(f"从配置文件加载人物列表成功，共 {len(CHARACTERS)} 个人物")
            
            self.log("配置文件加载成功")
        except Exception as e:
            QMessageBox.critical(self, "错误", f"加载配置文件失败: {str(e)}")
    
    def save_config(self):
        """保存配置"""
        # 更新全局配置
        CONFIG["ALIYUN_ACCESS_KEY_ID"] = self.access_key_id.text()
        CONFIG["ALIYUN_ACCESS_KEY_SECRET"] = self.access_key_secret.text()
        CONFIG["ALIYUN_APP_KEY"] = self.app_key.text()
        CONFIG["ALIYUN_REGION_ID"] = self.region_id.currentText()
        
        # 保存人物配置
        config_data = CONFIG.copy()
        config_data["CHARACTERS"] = CHARACTERS
        
        # 保存到文件
        config_file = "config.json"
        with open(config_file, "w", encoding="utf-8") as f:
            json.dump(config_data, f, indent=2, ensure_ascii=False)
        
        self.log(f"配置保存成功，文件: {config_file}")
        self.log(f"已保存 {len(CHARACTERS)} 个人物配置")
        QMessageBox.information(self, "成功", "配置保存成功")
    
    def test_connection(self):
        """测试连接"""
        # 更新全局配置
        CONFIG["ALIYUN_ACCESS_KEY_ID"] = self.access_key_id.text()
        CONFIG["ALIYUN_ACCESS_KEY_SECRET"] = self.access_key_secret.text()
        CONFIG["ALIYUN_APP_KEY"] = self.app_key.text()
        CONFIG["ALIYUN_REGION_ID"] = self.region_id.currentText()
        
        # 验证配置
        if not all([
            CONFIG["ALIYUN_ACCESS_KEY_ID"],
            CONFIG["ALIYUN_ACCESS_KEY_SECRET"],
            CONFIG["ALIYUN_APP_KEY"]
        ]):
            QMessageBox.warning(self, "警告", "请填写完整的配置信息")
            return
        
        # 测试获取token
        self.log("正在测试连接...")
        
        # 创建工作线程
        self.worker = WorkerThread("get_voices")
        self.worker.signal.connect(self.handle_worker_signal)
        self.worker.finished.connect(self.worker_thread_finished)
        self.worker.start()
    
    def get_voices(self):
        """获取音色列表"""
        # 验证配置
        if not all([
            CONFIG["ALIYUN_ACCESS_KEY_ID"],
            CONFIG["ALIYUN_ACCESS_KEY_SECRET"],
            CONFIG["ALIYUN_APP_KEY"]
        ]):
            QMessageBox.warning(self, "警告", "请先填写并保存阿里云配置")
            return
        
        self.log("正在获取音色列表...")
        
        # 创建工作线程
        self.worker = WorkerThread("get_voices")
        self.worker.signal.connect(self.handle_worker_signal)
        self.worker.finished.connect(self.worker_thread_finished)
        self.worker.start()
    
    def add_character(self):
        """添加人物"""
        name = self.character_name.text().strip()
        voice = self.character_voice.currentText()
        
        if not name:
            QMessageBox.warning(self, "警告", "请输入人物名称")
            return
        
        if not voice:
            QMessageBox.warning(self, "警告", "请选择音色")
            return
        
        # 检查是否已存在
        for char_info in CHARACTERS:
            if char_info.get("name") == name:
                QMessageBox.warning(self, "警告", "人物已存在")
                return
        
        # 添加人物
        CHARACTERS.append({
            "name": name,
            "voice": voice
        })
        
        # 更新人物列表
        self.update_character_list()
        
        # 清空输入
        self.character_name.clear()
        
        self.log(f"添加人物成功: {name} , 音色: {voice}")
    
    def delete_character(self):
        """删除选中人物"""
        selected_items = self.character_list.selectedItems()
        if not selected_items:
            QMessageBox.warning(self, "警告", "请选择要删除的人物")
            return
        
        for item in selected_items:
            character_name = item.text().split(" ")[0]
            # 从列表中删除
            for i, char_info in enumerate(CHARACTERS):
                if char_info.get("name") == character_name:
                    CHARACTERS.pop(i)
                    break
        
        # 更新人物列表
        self.update_character_list()
        
        self.log("删除人物成功")
    
    def update_character_list(self):
        """更新人物列表"""
        self.character_list.clear()
        for char_info in CHARACTERS:
            name = char_info.get("name")
            voice = char_info.get("voice")
            item = QListWidgetItem(f"{name} , 音色: {voice}")
            self.character_list.addItem(item)
    
    def browse_ppt_file(self):
        """浏览PPT文件"""
        file_path, _ = QFileDialog.getOpenFileName(self, "选择PPT文件", "", "PPT Files (*.pptx)")
        if file_path:
            self.ppt_file_path.setText(file_path)
            
            # 自动设置页面范围
            try:
                prs = Presentation(file_path)
                total_slides = len(prs.slides)
                self.end_page.setMaximum(total_slides)
                self.end_page.setValue(min(4, total_slides))
            except Exception:
                pass
    
    def process_ppt(self):
        """处理PPT"""
        ppt_file = self.ppt_file_path.text()
        start_page = self.start_page.value()
        end_page = self.end_page.value()
        
        if not ppt_file or not os.path.exists(ppt_file):
            QMessageBox.warning(self, "警告", "请选择有效的PPT文件")
            return
        
        if start_page > end_page:
            QMessageBox.warning(self, "警告", "起始页不能大于结束页")
            return
        
        # 验证配置
        if not all([
            CONFIG["ALIYUN_ACCESS_KEY_ID"],
            CONFIG["ALIYUN_ACCESS_KEY_SECRET"],
            CONFIG["ALIYUN_APP_KEY"]
        ]):
            QMessageBox.warning(self, "警告", "请先填写并保存阿里云配置")
            return
        
        self.log(f"开始处理PPT: {os.path.basename(ppt_file)}")
        self.log(f"处理页面范围: {start_page} - {end_page}")
        
        # 创建工作线程
        self.worker = WorkerThread("process_ppt",
                                   ppt_file=ppt_file,
                                   start_page=start_page,
                                   end_page=end_page)
        self.worker.signal.connect(self.handle_worker_signal)
        self.worker.finished.connect(self.worker_thread_finished)
        self.worker.start()
    
    def handle_worker_signal(self, signal_type, message):
        """处理工作线程信号"""
        if signal_type == "success":
            self.log(f"成功: {message}")
        elif signal_type == "error":
            self.log(f"错误: {message}")
        elif signal_type == "info":
            self.log(f"信息: {message}")
        
        # 更新音色列表
        if signal_type == "success" and "从Excel文件加载音色列表成功" in message:
            self.update_voice_list()
    
    def update_voice_list(self):
        """更新音色列表"""
        self.voice_list.clear()
        self.character_voice.clear()
        
        self.log(f"正在更新音色列表，共 {len(VOICES)} 个音色")
        
        for voice_info in VOICES:
            name = voice_info.get("name")
            description = voice_info.get("description")
            
            # 添加到音色列表
            item = QListWidgetItem(f"{name} , 描述: {description}")
            self.voice_list.addItem(item)
            
            # 添加到下拉框
            self.character_voice.addItem(name)
        
        self.log("音色列表更新完成")
    
    def test_selected_voice(self):
        """测试选中的音色"""
        # 验证配置
        if not all([
            CONFIG["ALIYUN_ACCESS_KEY_ID"],
            CONFIG["ALIYUN_ACCESS_KEY_SECRET"],
            CONFIG["ALIYUN_APP_KEY"]
        ]):
            QMessageBox.warning(self, "警告", "请先填写并保存阿里云配置")
            return
        
        # 获取选中的音色
        selected_items = self.voice_list.selectedItems()
        if not selected_items:
            QMessageBox.warning(self, "警告", "请选择要测试的音色")
            return
        
        # 提取音色名称
        voice_item = selected_items[0].text()
        voice_name = voice_item.split(" ")[0]
        
        self.log(f"正在测试音色: {voice_name}")
        
        # 创建测试文本
        test_text = "这是一段测试文本，用于测试音色效果。"
        
        # 创建测试音频文件
        test_audio_file = os.path.join("audio_files", f"test_{voice_name}.mp3")
        if not os.path.exists("audio_files"):
            os.makedirs("audio_files")
        
        # 创建并启动测试线程
        # 使用类成员变量保存线程对象，避免被垃圾回收
        self.test_thread = TestVoiceThread(voice_name, test_text, test_audio_file)
        self.test_thread.signal.connect(self.handle_worker_signal)
        self.test_thread.finished.connect(lambda: self.test_thread_finished(voice_name))
        self.test_thread.start()

    def test_thread_finished(self, voice_name):
        """测试线程完成后的处理"""
        # 线程完成后，可以进行一些清理工作
        self.log(f"音色测试线程已完成: {voice_name}")
        # 释放线程对象
        self.test_thread = None
    
    def worker_thread_finished(self):
        """工作线程完成后的处理"""
        # 线程完成后，可以进行一些清理工作
        # self.log("工作线程已完成")
        # 释放线程对象
        self.worker = None
    
    def log(self, message):
        """记录日志"""
        timestamp = time.strftime("%Y-%m-%d %H:%M:%S", time.localtime())
        self.status_text.append(f"[{timestamp}] {message}")
        self.status_text.verticalScrollBar().setValue(self.status_text.verticalScrollBar().maximum())

if __name__ == "__main__":
    app = QApplication(sys.argv)
    window = MainWindow()
    window.show()
    sys.exit(app.exec_())